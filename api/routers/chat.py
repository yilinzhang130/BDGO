"""Chat streaming endpoint — MiniMax API with tool use + document attachments."""

import asyncio
import json, logging, uuid
from pathlib import Path
import httpx
from fastapi import APIRouter, Depends
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from db import query, query_one, count
from config import BP_DIR
from database import transaction
from auth import get_current_user
from models import resolve_model, ModelSpec
from field_policy import strip_hidden, HIDDEN_FIELDS
import credits as credits_mod
import planner as planner_mod

logger = logging.getLogger(__name__)
router = APIRouter()


# ---------------------------------------------------------------------------
# Postgres-backed session helpers
# ---------------------------------------------------------------------------

def _ensure_session(session_id: str, user_id: str) -> None:
    """Create the session row if it doesn't exist yet."""
    with transaction() as cur:
        cur.execute("SELECT id FROM sessions WHERE id = %s", (session_id,))
        if not cur.fetchone():
            cur.execute(
                "INSERT INTO sessions (id, user_id, title) VALUES (%s, %s, %s)",
                (session_id, user_id, "New Chat"),
            )


# Context-compaction tuning knobs
COMPACT_KEEP_VERBATIM = 4            # last N turns (user+assistant pairs) kept as-is
COMPACT_TOKEN_BUDGET = 40_000        # target tokens after layer 1; triggers layer 2 if exceeded
COMPACT_CHARS_PER_TOKEN = 2.5        # rough estimator for mixed CJK+ASCII


def _estimate_tokens(history: list[dict]) -> int:
    """Cheap char-based token estimate for a history list."""
    total_chars = 0
    for m in history:
        c = m.get("content")
        if isinstance(c, list):
            for b in c:
                if not isinstance(b, dict):
                    continue
                # text / tool_use input / tool_result content all count
                t = b.get("text")
                if isinstance(t, str):
                    total_chars += len(t)
                tr = b.get("content")
                if isinstance(tr, list):
                    for sub in tr:
                        if isinstance(sub, dict):
                            total_chars += len(sub.get("text", "") or "")
                elif isinstance(tr, str):
                    total_chars += len(tr)
                inp = b.get("input")
                if isinstance(inp, dict):
                    total_chars += len(json.dumps(inp, ensure_ascii=False))
        elif isinstance(c, str):
            total_chars += len(c)
    return int(total_chars / COMPACT_CHARS_PER_TOKEN)


def _strip_tool_blocks_from_old(history: list[dict], keep_last_n: int) -> list[dict]:
    """Layer 1: strip non-text content blocks from messages older than the last
    `keep_last_n` user turns. Preserves tool/result pairs in the recent zone so
    the LLM can still reason over them.
    """
    user_idx = [i for i, m in enumerate(history) if m.get("role") == "user"]
    if len(user_idx) <= keep_last_n:
        return history  # nothing to strip

    boundary = user_idx[-keep_last_n]  # messages before this are "old"

    result: list[dict] = []
    for i, m in enumerate(history):
        if i >= boundary:
            result.append(m)
            continue

        content = m.get("content")
        if isinstance(content, list):
            # Keep only text blocks in old zone
            text_blocks = [
                b for b in content
                if isinstance(b, dict) and b.get("type") == "text" and b.get("text", "").strip()
            ]
            if text_blocks:
                result.append({"role": m["role"], "content": text_blocks})
            # else: message was only tool_use/tool_result/empty — drop it
        elif isinstance(content, str) and content.strip():
            result.append(m)
        # else: empty — drop

    return result


def _load_history(session_id: str) -> list[dict]:
    """Load the last 20 messages for a session from Postgres.

    Messages are stored with role + content.  For assistant messages the
    content column holds a JSON-encoded list of content blocks (text /
    tool_use).  For user messages with tool_results, the content column
    holds the JSON-encoded list.

    Plan-card placeholder messages (assistant with empty content + a plan
    stored in tools_json) are skipped — they're UI-only artifacts; the LLM
    never needs to see "the user saw a plan card at this point".
    """
    with transaction() as cur:
        cur.execute(
            "SELECT role, content, tools_json FROM messages "
            "WHERE session_id = %s ORDER BY created_at DESC LIMIT 100",
            (session_id,),
        )
        rows = cur.fetchall()

    # Rows come back newest-first; reverse to chronological order.
    rows.reverse()

    history: list[dict] = []
    for r in rows:
        content = r["content"]
        tools_json = r.get("tools_json")

        # Skip plan placeholders: empty-content assistant with a plan blob.
        if r["role"] == "assistant" and not (content or "").strip() and tools_json:
            try:
                parsed_tools = json.loads(tools_json)
                if isinstance(parsed_tools, dict) and parsed_tools.get("plan"):
                    continue
            except (json.JSONDecodeError, TypeError):
                pass

        # Try parsing JSON content (structured content blocks / tool results)
        try:
            parsed = json.loads(content)
            if isinstance(parsed, list):
                history.append({"role": r["role"], "content": parsed})
                continue
        except (json.JSONDecodeError, TypeError):
            pass
        history.append({"role": r["role"], "content": content})

    # No hard cap here — _compact_if_needed decides what to drop / summarize
    # based on token budget, not raw count.
    return history


def _save_message(session_id: str, role: str, content, tools_json: str | None = None,
                  attachments_json: str | None = None) -> None:
    """Persist a single message to Postgres."""
    if isinstance(content, (list, dict)):
        content_str = json.dumps(content, ensure_ascii=False, default=str)
    else:
        content_str = str(content) if content else ""

    try:
        with transaction() as cur:
            cur.execute(
                "INSERT INTO messages (id, session_id, role, content, tools_json, attachments_json) "
                "VALUES (%s, %s, %s, %s, %s, %s)",
                (uuid.uuid4().hex[:12], session_id, role, content_str, tools_json, attachments_json),
            )
            cur.execute(
                "UPDATE sessions SET updated_at = NOW() WHERE id = %s",
                (session_id,),
            )
    except Exception:
        logger.exception("Failed to save message for session %s", session_id)


def _get_session_brief(session_id: str) -> tuple[str | None, str | None]:
    """Return (brief_text, brief_ts_iso) for a session, or (None, None)."""
    try:
        with transaction() as cur:
            cur.execute(
                "SELECT brief, brief_ts FROM sessions WHERE id = %s",
                (session_id,),
            )
            row = cur.fetchone()
    except Exception:
        logger.exception("Failed to load session brief")
        return (None, None)
    if not row:
        return (None, None)
    brief = row.get("brief")
    ts = row.get("brief_ts")
    return (brief, ts.isoformat() if ts else None)


def _save_session_brief(session_id: str, brief: str, brief_ts: str | None) -> None:
    try:
        with transaction() as cur:
            cur.execute(
                "UPDATE sessions SET brief = %s, brief_ts = COALESCE(%s::timestamp, brief_ts) "
                "WHERE id = %s",
                (brief, brief_ts, session_id),
            )
    except Exception:
        logger.exception("Failed to save session brief")


async def _compact_if_needed(
    session_id: str,
    history: list[dict],
    model: "ModelSpec",
) -> list[dict]:
    """Apply layer-1 (strip old tool blocks) + layer-2 (LLM summary) if
    the estimated token count exceeds the budget. Returns the compacted
    history ready to pass to the LLM. Cached summary is reused across
    turns so we only pay the summarization cost once per high-watermark.
    """
    # Layer 1: always applied — cheap and reduces noise even when not over budget
    compacted = _strip_tool_blocks_from_old(history, COMPACT_KEEP_VERBATIM)

    # Fast path: within budget, just return layer-1 result
    est = _estimate_tokens(compacted)
    if est <= COMPACT_TOKEN_BUDGET:
        return compacted

    # Layer 2: prepare the old turns (everything before the keep-verbatim zone)
    # plus any existing brief, ask the summarizer LLM to produce a new brief.
    user_idx = [i for i, m in enumerate(compacted) if m.get("role") == "user"]
    if len(user_idx) <= COMPACT_KEEP_VERBATIM:
        return compacted  # can't compact further

    boundary = user_idx[-COMPACT_KEEP_VERBATIM]
    old_turns = compacted[:boundary]
    recent = compacted[boundary:]

    existing_brief, _ = _get_session_brief(session_id)

    # Cheap cache-hit path: if a brief already exists and substituting it for
    # the old turns would be within budget, reuse it — no LLM call needed.
    if existing_brief:
        estimated_with_cached = _estimate_tokens([
            {"role": "user", "content": existing_brief}
        ] + recent)
        if estimated_with_cached <= COMPACT_TOKEN_BUDGET:
            brief_turn = {
                "role": "user",
                "content": f"[会话早前内容的摘要，供你参考]\n{existing_brief}\n\n[以下是最近几轮对话的原文]",
            }
            return [brief_turn] + recent

    # Otherwise regenerate brief incorporating old_turns + prior brief.
    new_brief = await planner_mod.summarize_history(old_turns, existing_brief, model)
    if not new_brief:
        # Summarizer failed — fall back to layer-1 only
        return compacted

    # Cache the new brief (idempotent across turns)
    _save_session_brief(session_id, new_brief, None)

    # Compose final history: a single assistant turn carrying the brief,
    # then recent verbatim turns. Wrapping as assistant text works with
    # any Anthropic-compat provider and doesn't require extra role support.
    brief_turn = {
        "role": "user",
        "content": f"[会话早前内容的摘要，供你参考]\n{new_brief}\n\n[以下是最近几轮对话的原文]",
    }
    return [brief_turn] + recent



    """Upsert context entities extracted from tool results."""
    if not entities:
        return
    try:
        with transaction() as cur:
            params = [
                (e.get("id"), session_id, e.get("entity_type", ""),
                 e.get("title", ""), e.get("subtitle"),
                 json.dumps(e.get("fields", []), ensure_ascii=False) if e.get("fields") else None,
                 e.get("href"))
                for e in entities
            ]
            cur.executemany(
                """INSERT INTO context_entities (id, session_id, entity_type, title, subtitle, fields_json, href)
                   VALUES (%s, %s, %s, %s, %s, %s, %s)
                   ON CONFLICT (id) DO UPDATE SET
                       entity_type = EXCLUDED.entity_type,
                       title = EXCLUDED.title,
                       subtitle = EXCLUDED.subtitle,
                       fields_json = EXCLUDED.fields_json,
                       href = EXCLUDED.href,
                       added_at = NOW()""",
                params,
            )
    except Exception:
        logger.exception("Failed to save entities for session %s", session_id)

SYSTEM_PROMPT = """你是 BD Go，生物医药BD领域的资深智能助手。你不是通用AI，你是一位有10年BD经验的VP级同事。

## 你的工具箱

**BDGO 数据库**（只读）：
- search_companies / get_company — 公司检索
- search_assets / get_asset — 管线资产
- search_clinical — 临床试验（44,000+条）
- search_deals — BD交易（7,000+条）
- search_patents — 专利
- get_buyer_profile — MNC买方画像
- count_by — 聚合统计
- search_global — 跨表搜索

**临床指南库**（79条指南, 611条推荐, 379条biomarker）：
- query_treatment_guidelines — 治疗推荐
- query_biomarker — 生物标志物
- list_guidelines — 指南列表

**报告生成**（返回Word下载链接，异步任务约2-3分钟）：
- generate_buyer_profile — MNC买方画像报告
- research_disease — 赛道竞争格局报告（8章，含立项评分矩阵）
- analyze_commercial — 资产商业化评估（患者漏斗/TAM/定价/Revenue Forecast三情景）
- analyze_ip — 专利景观简报（6章）
- analyze_target — 靶点竞争雷达
- analyze_paper — 文献综述
- generate_guidelines_report — 临床指南简报

## 核心行为原则

### 1. 永远主动出击，不要被动等待

你收到任何问题后，第一件事是判断意图，第二件事是**立即并行调用多个工具**。

❌ 错误："请问您想了解哪方面？"
✅ 正确：直接调用3-4个工具，收集数据后给出完整分析

### 2. 附件 = 自动触发深度分析

当用户上传文件（BP、报告、文档），**不要只是总结文件内容**。你必须：

**Step 1**: 从文件中提取关键实体（公司名、资产名、靶点、适应症、临床阶段）
**Step 2**: 用提取的实体**立即查 BDGO 数据库**：
  - search_companies → 这家公司在数据库里有没有？什么背景？
  - search_assets → 他们的管线我们有记录吗？竞品情况？
  - search_clinical → 临床数据怎样？
  - search_deals → 有过BD交易吗？可比交易？
  - query_treatment_guidelines → 目标适应症的治疗格局如何？
**Step 3**: 交叉验证文件声称 vs BDGO数据库数据，找出差异和亮点
**Step 4**: 输出结构化分析：
  ```
  📋 文件概要：[1-2句]
  🏢 公司画像：[BDGO数据库 vs 文件声称]
  💊 管线评估：[临床阶段、竞品对比、差异化]
  🏥 治疗格局：[指南推荐、unmet need]
  💰 交易参考：[可比交易、估值区间]
  ⚡ BD建议：[追/观望/放弃 + 理由]
  ```

### 3. 智能意图路由

| 用户说的 | 你要做的 |
|---------|---------|
| "分析这个公司/这个BP" | 全套：公司+资产+临床+交易+指南，5步分析 |
| "XX赛道怎么样" | search_assets(disease) + search_clinical + query_treatment_guidelines + search_deals + count_by |
| "XX靶点值得追吗" | search_assets(target) + search_clinical + query_biomarker + search_deals，给出go/no-go建议 |
| "帮我看看XX公司" | get_company + search_assets + search_deals + get_buyer_profile |
| "XX和YY比怎么样" | 两边都查，做head-to-head对比表格 |
| "最近有什么deal" | search_deals(sort_by=date) + 分析趋势 |
| "XX赛道报告/竞争格局" | research_disease(disease=XX) — 生成8章Word报告 |
| "XX商业化/peak sales/能卖多少" | analyze_commercial(asset_name=XX) — 患者漏斗+Revenue Forecast报告 |
| "XX专利格局/IP landscape" | analyze_ip(query=XX) — 专利景观简报 |
| "帮我做个报告" | 用 research_disease / analyze_commercial / analyze_ip 等工具，告知用户去Reports页面下载 |
| "帮我加关注/收藏/加入关注" | add_to_watchlist(entity_type, entity_key) — 立即执行，告知用户已添加成功 |
| 简单查数据 | 直接一个工具，快速回答 |

### 4. 每轮输出要有结构

不要给一堆散乱的文字。用以下格式之一：

**快速查询** → 表格 + 1-2句分析
**深度研究** → 带emoji分区的结构化报告（见Step 4）
**对比分析** → 对比表格 + 优劣势总结
**数据统计** → 数字 + 趋势判断

### 5. 工具调用规则

- **并行调用**：如果多个工具之间无依赖，在同一轮一起调用（如search_assets和search_deals可以同时调用）
- **Coverage Check**：每轮工具返回后检查，如有关键信息缺失，补充调用（最多3轮）
- **必须用工具**：有数据库数据可查时，绝不编造答案
- **数据引用**：注明来源（"来源：BDGO资产库"、"来源：NCCN 2025指南"）

### 6. 语言和沟通风格

- 中文为主，专业术语保留英文（MOA、PFS、ORR、DCR）
- 像资深BD同事，不像AI — 有观点、有判断、敢说"这个不值得追"
- 简洁直接，不重复用户问题，不说"好的，让我来帮你"
"""

# ── Tool definitions ─────────────────────────────────────────

TOOLS = [
    {
        "name": "search_companies",
        "description": "Search companies in CRM by name/country/type/disease. Returns top matches.",
        "input_schema": {
            "type": "object",
            "properties": {
                "q": {"type": "string", "description": "Search keyword (matches company name)"},
                "country": {"type": "string", "description": "e.g. USA, China, Japan"},
                "type": {"type": "string", "description": "e.g. Biotech(USA), 海外药企"},
                "disease": {"type": "string", "description": "e.g. Oncology, Immunology"},
                "limit": {"type": "integer", "default": 10},
            },
        },
    },
    {
        "name": "get_company",
        "description": "Get full details of one company by exact name.",
        "input_schema": {
            "type": "object",
            "properties": {"name": {"type": "string"}},
            "required": ["name"],
        },
    },
    {
        "name": "search_assets",
        "description": "Search pipeline assets by company, phase, disease, target.",
        "input_schema": {
            "type": "object",
            "properties": {
                "q": {"type": "string"},
                "company": {"type": "string"},
                "phase": {"type": "string", "description": "e.g. Phase 1, Phase 2, Phase 3, Commercial"},
                "disease": {"type": "string"},
                "target": {"type": "string"},
                "limit": {"type": "integer", "default": 10},
            },
        },
    },
    {
        "name": "get_asset",
        "description": "Get full details of one asset by name and company (composite key).",
        "input_schema": {
            "type": "object",
            "properties": {
                "name": {"type": "string"},
                "company": {"type": "string"},
            },
            "required": ["name", "company"],
        },
    },
    {
        "name": "search_clinical",
        "description": "Search clinical trials by company/asset/phase/indication.",
        "input_schema": {
            "type": "object",
            "properties": {
                "q": {"type": "string"},
                "company": {"type": "string"},
                "asset": {"type": "string"},
                "phase": {"type": "string"},
                "limit": {"type": "integer", "default": 10},
            },
        },
    },
    {
        "name": "search_deals",
        "description": "Search BD deals/transactions by buyer, seller, or deal type.",
        "input_schema": {
            "type": "object",
            "properties": {
                "q": {"type": "string"},
                "buyer": {"type": "string"},
                "seller": {"type": "string"},
                "type": {"type": "string", "description": "e.g. M&A, License, Collaboration"},
                "sort_by": {"type": "string", "enum": ["date", "upfront", "total"], "description": "Sort field"},
                "limit": {"type": "integer", "default": 10},
            },
        },
    },
    {
        "name": "search_patents",
        "description": "Search patents by company, asset, or status.",
        "input_schema": {
            "type": "object",
            "properties": {
                "q": {"type": "string"},
                "company": {"type": "string"},
                "status": {"type": "string", "description": "有效 / 已过期"},
                "limit": {"type": "integer", "default": 10},
            },
        },
    },
    {
        "name": "get_buyer_profile",
        "description": "Get MNC buyer profile (DNA summary, BD theses, commercial capabilities) for a company.",
        "input_schema": {
            "type": "object",
            "properties": {"company": {"type": "string"}},
            "required": ["company"],
        },
    },
    {
        "name": "count_by",
        "description": "Count records grouped by a column. Use for aggregation questions like '按阶段统计资产数量'.",
        "input_schema": {
            "type": "object",
            "properties": {
                "table": {"type": "string", "enum": ["公司", "资产", "临床", "交易", "IP"]},
                "group_by": {"type": "string", "description": "Column name to group by, e.g. 临床阶段, 疾病领域"},
            },
            "required": ["table", "group_by"],
        },
    },
    {
        "name": "search_global",
        "description": "Fuzzy search across all tables (companies + assets + clinical + deals + patents).",
        "input_schema": {
            "type": "object",
            "properties": {
                "q": {"type": "string"},
                "limit": {"type": "integer", "default": 5},
            },
            "required": ["q"],
        },
    },
    {
        "name": "add_to_watchlist",
        "description": "Add a company or asset to the current user's watchlist. Use this when the user says '加关注', '加入关注', '收藏', or any similar intent. entity_type must be 'company' or 'asset'.",
        "input_schema": {
            "type": "object",
            "properties": {
                "entity_type": {"type": "string", "enum": ["company", "asset", "disease", "target", "incubator"], "description": "Type of entity"},
                "entity_key": {"type": "string", "description": "The name/key of the entity, e.g. the company name or asset name"},
                "notes": {"type": "string", "description": "Optional notes about why this is being watched"},
            },
            "required": ["entity_type", "entity_key"],
        },
    },
    # ── Clinical Guidelines tools (from guidelines.db) ──
    {
        "name": "query_treatment_guidelines",
        "description": "Query clinical treatment guidelines by disease. Returns recommended drugs by treatment line with evidence grade. Covers NCCN, CSCO, ESMO, and 20+ guideline sources across 74 diseases.",
        "input_schema": {
            "type": "object",
            "properties": {
                "disease": {"type": "string", "description": "Disease name in Chinese or English, e.g. '非小细胞肺癌', 'NSCLC', 'breast cancer', '乳腺癌'"},
                "line": {"type": "string", "description": "Treatment line filter, e.g. '一线', '二线', '三线'. Omit for all lines."},
                "source": {"type": "string", "description": "Guideline source filter, e.g. 'NCCN', 'CSCO', 'ESMO'. Omit for all sources."},
                "limit": {"type": "integer", "default": 20},
            },
            "required": ["disease"],
        },
    },
    {
        "name": "query_biomarker",
        "description": "Query biomarker information for a disease — testing methods, positive thresholds, clinical significance. Covers 379 biomarker records from major guidelines.",
        "input_schema": {
            "type": "object",
            "properties": {
                "disease": {"type": "string", "description": "Disease name, e.g. '非小细胞肺癌', 'breast cancer'"},
                "biomarker": {"type": "string", "description": "Optional: specific biomarker name, e.g. 'PD-L1', 'EGFR', 'HER2'. Omit for all biomarkers."},
            },
            "required": ["disease"],
        },
    },
    {
        "name": "list_guidelines",
        "description": "List available clinical guidelines by disease or source. Shows guideline name, source (NCCN/CSCO/ESMO), version, and year. Use this when the user asks 'what guidelines are available for X?'",
        "input_schema": {
            "type": "object",
            "properties": {
                "disease": {"type": "string", "description": "Disease name filter. Omit for all diseases."},
                "source": {"type": "string", "description": "Source filter (NCCN/CSCO/ESMO/etc). Omit for all sources."},
            },
        },
    },
]


# ── Tool implementations ─────────────────────────────────────

def _tool_search_companies(q="", country="", type="", disease="", limit=10):
    conds, params = [], []
    if q:
        conds.append('"客户名称" LIKE ?')
        params.append(f"%{q}%")
    if country:
        conds.append('"所处国家" = ?')
        params.append(country)
    if type:
        conds.append('"客户类型" = ?')
        params.append(type)
    if disease:
        conds.append('"疾病领域" LIKE ?')
        params.append(f"%{disease}%")
    where = " AND ".join(conds) if conds else "1=1"
    sql = f'SELECT "客户名称","客户类型","所处国家","疾病领域","核心产品的阶段","BD跟进优先级","公司质量评分" FROM "公司" WHERE {where} LIMIT ?'
    params.append(min(limit, 30))
    return query(sql, tuple(params))


def _tool_get_company(name):
    return query_one('SELECT * FROM "公司" WHERE "客户名称" = ?', (name,))


def _tool_search_assets(q="", company="", phase="", disease="", target="", limit=10):
    conds, params = [], []
    if q:
        conds.append('("资产名称" LIKE ? OR "资产代号" LIKE ?)')
        params.extend([f"%{q}%", f"%{q}%"])
    if company:
        conds.append('"所属客户" LIKE ?')
        params.append(f"%{company}%")
    if phase:
        conds.append('"临床阶段" = ?')
        params.append(phase)
    if disease:
        conds.append('"疾病领域" LIKE ?')
        params.append(f"%{disease}%")
    if target:
        conds.append('"靶点" LIKE ?')
        params.append(f"%{target}%")
    where = " AND ".join(conds) if conds else "1=1"
    sql = f'SELECT "资产名称","所属客户","临床阶段","疾病领域","适应症","靶点","作用机制(MOA)","Q总分" FROM "资产" WHERE {where} LIMIT ?'
    params.append(min(limit, 30))
    return query(sql, tuple(params))


def _tool_get_asset(name, company):
    return query_one('SELECT * FROM "资产" WHERE "资产名称" = ? AND "所属客户" = ?', (name, company))


def _tool_search_clinical(q="", company="", asset="", phase="", limit=10):
    conds, params = [], []
    if q:
        conds.append('("试验ID" LIKE ? OR "适应症" LIKE ?)')
        params.extend([f"%{q}%", f"%{q}%"])
    if company:
        conds.append('"公司名称" LIKE ?')
        params.append(f"%{company}%")
    if asset:
        conds.append('"资产名称" LIKE ?')
        params.append(f"%{asset}%")
    if phase:
        conds.append('"临床期次" = ?')
        params.append(phase)
    where = " AND ".join(conds) if conds else "1=1"
    sql = f'SELECT "记录ID","试验ID","公司名称","资产名称","适应症","临床期次","主要终点名称","结果判定","数据状态" FROM "临床" WHERE {where} LIMIT ?'
    params.append(min(limit, 30))
    return query(sql, tuple(params))


def _tool_search_deals(q="", buyer="", seller="", type="", sort_by="date", limit=10):
    conds, params = [], []
    if q:
        conds.append('"交易名称" LIKE ?')
        params.append(f"%{q}%")
    if buyer:
        conds.append('"买方公司" LIKE ?')
        params.append(f"%{buyer}%")
    if seller:
        conds.append('"卖方/合作方" LIKE ?')
        params.append(f"%{seller}%")
    if type:
        conds.append('"交易类型" LIKE ?')
        params.append(f"%{type}%")
    where = " AND ".join(conds) if conds else "1=1"
    order_col = {"date": "宣布日期", "upfront": "首付款($M)", "total": "交易总额($M)"}.get(sort_by, "宣布日期")
    sql = f'SELECT "交易名称","交易类型","买方公司","卖方/合作方","资产名称","首付款($M)","交易总额($M)","宣布日期","战略评分" FROM "交易" WHERE {where} ORDER BY "{order_col}" DESC LIMIT ?'
    params.append(min(limit, 30))
    return query(sql, tuple(params))


def _tool_search_patents(q="", company="", status="", limit=10):
    conds, params = [], []
    if q:
        conds.append('"专利号" LIKE ?')
        params.append(f"%{q}%")
    if company:
        conds.append('"关联公司" LIKE ?')
        params.append(f"%{company}%")
    if status:
        conds.append('"状态" = ?')
        params.append(status)
    where = " AND ".join(conds) if conds else "1=1"
    sql = f'SELECT "专利号","专利持有人","关联公司","关联资产","到期日","状态","管辖区" FROM "IP" WHERE {where} LIMIT ?'
    params.append(min(limit, 30))
    return query(sql, tuple(params))


def _tool_get_buyer_profile(company):
    return query_one('SELECT * FROM "MNC画像" WHERE "company_name" = ?', (company,))


def _tool_count_by(table, group_by):
    table_map = {"公司": "公司", "资产": "资产", "临床": "临床", "交易": "交易", "IP": "IP"}
    physical = table_map.get(table, table)
    if table not in table_map:
        return {"error": f"Invalid table: {table}"}
    # Validate column name against actual table columns to prevent SQL injection
    try:
        cols = [r["name"] for r in query(f'PRAGMA table_info("{physical}")')]
    except Exception:
        cols = []
    if group_by not in cols:
        return {"error": f"Invalid column '{group_by}' for table '{table}'"}
    sql = f'SELECT "{group_by}" as value, COUNT(*) as count FROM "{physical}" WHERE "{group_by}" IS NOT NULL AND "{group_by}" != \'\' GROUP BY "{group_by}" ORDER BY count DESC LIMIT 50'
    try:
        return query(sql)
    except Exception as e:
        return {"error": str(e)}


def _tool_search_global(q, limit=5):
    results = {}
    results["companies"] = _tool_search_companies(q=q, limit=limit)
    results["assets"] = _tool_search_assets(q=q, limit=limit)
    results["deals"] = _tool_search_deals(q=q, limit=limit)
    results["clinical"] = _tool_search_clinical(q=q, limit=limit)
    return results


def _tool_add_to_watchlist(entity_type: str, entity_key: str, notes: str = "", _user_id: str = ""):
    """Add entity to the current user's watchlist (writes to PostgreSQL)."""
    if not _user_id:
        return {"success": False, "error": "用户未登录，无法添加关注"}
    try:
        with transaction() as cur:
            cur.execute(
                """INSERT INTO user_watchlists (entity_type, entity_key, user_id, notes)
                   VALUES (%s, %s, %s, %s)
                   ON CONFLICT (user_id, entity_type, entity_key)
                   DO UPDATE SET notes = COALESCE(NULLIF(EXCLUDED.notes, ''), user_watchlists.notes)
                   RETURNING id""",
                (entity_type, entity_key, _user_id, notes or None),
            )
            row = cur.fetchone()
        return {
            "success": True,
            "id": row["id"],
            "entity_type": entity_type,
            "entity_key": entity_key,
            "message": f"✅ 已将 **{entity_key}** 添加到关注列表！你可以在 [关注列表](/watchlist) 页面查看。",
        }
    except Exception as e:
        logger.exception("add_to_watchlist failed")
        return {"success": False, "error": str(e)}


# ── Clinical Guidelines tools (guidelines.db) ───────────────

import os
import sqlite3 as _sqlite3

from config import GUIDELINES_DB_PATH as _GUIDELINES_DB


def _guidelines_query(sql: str, params: tuple = ()) -> list[dict]:
    """Query guidelines.db (separate from CRM). Returns list of dicts."""
    if not os.path.exists(_GUIDELINES_DB):
        return [{"error": "Guidelines database not found"}]
    conn = _sqlite3.connect(f"file:{_GUIDELINES_DB}?mode=ro", uri=True)
    conn.row_factory = _sqlite3.Row
    try:
        rows = conn.execute(sql, params).fetchall()
        return [dict(r) for r in rows]
    finally:
        conn.close()


def _tool_query_treatment_guidelines(disease="", line="", source="", limit=20):
    conds = ['g."疾病" LIKE ?']
    params: list = [f"%{disease}%"]
    if line:
        conds.append('r."治疗线" LIKE ?')
        params.append(f"%{line}%")
    if source:
        conds.append('g."指南来源" LIKE ?')
        params.append(f"%{source}%")
    where = " AND ".join(conds)
    sql = f'''
        SELECT g."疾病", g."疾病亚型", g."指南来源", g."版本", g."发布年份",
               r."治疗线", r."推荐等级", r."证据级别", r."药物", r."药物类型",
               r."给药方案", r."适应条件", r."适应人群", r."疗效概述", r."不良反应", r."备注"
        FROM "推荐" r
        JOIN "指南" g ON r."指南ID" = g."记录ID"
        WHERE {where}
        ORDER BY g."疾病", r."治疗线", r."推荐等级"
        LIMIT ?
    '''
    params.append(min(limit, 50))
    return _guidelines_query(sql, tuple(params))


def _tool_query_biomarker(disease="", biomarker=""):
    conds = ['"疾病" LIKE ?']
    params: list = [f"%{disease}%"]
    if biomarker:
        conds.append('"标志物" LIKE ?')
        params.append(f"%{biomarker}%")
    where = " AND ".join(conds)
    sql = f'''
        SELECT "疾病", "标志物", "检测方法", "阳性阈值", "临床意义", "指南来源", "最近更新"
        FROM "生物标志物"
        WHERE {where}
        ORDER BY "疾病", "标志物"
        LIMIT 30
    '''
    return _guidelines_query(sql, tuple(params))


def _tool_list_guidelines(disease="", source=""):
    conds: list[str] = []
    params: list = []
    if disease:
        conds.append('"疾病" LIKE ?')
        params.append(f"%{disease}%")
    if source:
        conds.append('"指南来源" LIKE ?')
        params.append(f"%{source}%")
    where = " AND ".join(conds) if conds else "1=1"
    sql = f'''
        SELECT "疾病", "疾病亚型", "指南来源", "版本", "发布年份", "URL"
        FROM "指南"
        WHERE {where}
        ORDER BY "疾病", "指南来源"
        LIMIT 50
    '''
    return _guidelines_query(sql, tuple(params))


TOOL_IMPL = {
    "search_companies": _tool_search_companies,
    "get_company": _tool_get_company,
    "search_assets": _tool_search_assets,
    "get_asset": _tool_get_asset,
    "search_clinical": _tool_search_clinical,
    "search_deals": _tool_search_deals,
    "search_patents": _tool_search_patents,
    "get_buyer_profile": _tool_get_buyer_profile,
    "count_by": _tool_count_by,
    "search_global": _tool_search_global,
    "add_to_watchlist": _tool_add_to_watchlist,
    # Guidelines tools
    "query_treatment_guidelines": _tool_query_treatment_guidelines,
    "query_biomarker": _tool_query_biomarker,
    "list_guidelines": _tool_list_guidelines,
}


# ── Auto-register report services as Chat tools ─────────────
# Each registered ReportService becomes a MiniMax tool.
# - async service: tool kicks off a task, returns task_id + a user-facing message
# - sync service: tool runs inline and returns the markdown preview
#
# Adding a new report in services/__init__.py automatically exposes it here.
try:
    from services import REPORT_SERVICES
    from services.report_builder import create_task, execute_task, get_task
    import threading as _threading
    import json as _json_rep

    def _persist_report(task_id: str, slug: str, user_id: str) -> None:
        """Save completed report to report_history so it shows up in My Reports."""
        task = get_task(task_id)
        if not (task and task.get("status") == "completed"):
            return
        result = task.get("result") or {}
        title = result.get("meta", {}).get("title") or slug
        markdown_preview = (result.get("markdown") or "")[:2000]
        files_json = _json_rep.dumps(result.get("files", []), ensure_ascii=False, default=str)
        meta_json = _json_rep.dumps(result.get("meta", {}), ensure_ascii=False, default=str)
        try:
            from database import transaction
            with transaction() as cur:
                cur.execute(
                    "INSERT INTO report_history (user_id, task_id, slug, title, markdown_preview, files_json, meta_json) "
                    "VALUES (%s, %s, %s, %s, %s, %s, %s) ON CONFLICT (task_id) DO NOTHING",
                    (user_id, task_id, slug, title, markdown_preview, files_json, meta_json),
                )
        except Exception:
            logger.exception("Failed to persist report history for task %s", task_id)

    def _make_report_tool(_svc):
        def _impl(_user_id=None, **kwargs):
            task_id = create_task(_svc.slug, kwargs)
            if _svc.mode == "sync":
                execute_task(task_id, _svc, kwargs)
                t = get_task(task_id)
                if t["status"] == "completed":
                    if _user_id:
                        _persist_report(task_id, _svc.slug, _user_id)
                    result = t.get("result") or {}
                    return {
                        "task_id": task_id,
                        "status": "completed",
                        "markdown_preview": (result.get("markdown") or "")[:3000],
                        "files": result.get("files", []),
                        "message": f"Report generated. Download: {', '.join(f['download_url'] for f in result.get('files', []))}",
                    }
                return {
                    "task_id": task_id,
                    "status": t["status"],
                    "error": t.get("error"),
                }
            # Async: spawn thread, return immediately
            def _run():
                execute_task(task_id, _svc, kwargs)
                if _user_id:
                    _persist_report(task_id, _svc.slug, _user_id)
            thread = _threading.Thread(target=_run, daemon=True)
            thread.start()
            return {
                "task_id": task_id,
                "status": "queued",
                "estimated_seconds": _svc.estimated_seconds,
                "display_message": (
                    f"已开始 {_svc.display_name}（任务ID: {task_id}）。"
                    f"预计 {_svc.estimated_seconds} 秒完成。"
                    f"可在 Reports 页面查看进度与下载结果。"
                ),
            }
        return _impl

    for _svc in REPORT_SERVICES.values():
        TOOLS.append({
            "name": _svc.chat_tool_name,
            "description": _svc.chat_tool_description,
            "input_schema": _svc.chat_tool_input_schema,
        })
        TOOL_IMPL[_svc.chat_tool_name] = _make_report_tool(_svc)

    logger.info("Registered %d report services as Chat tools", len(REPORT_SERVICES))
except Exception as _e:
    logger.exception("Failed to register report services as Chat tools: %s", _e)


# Map each CRM tool to its source table so we can apply field visibility.
# Tools NOT in this map (watchlist / reports / skill calls) are not filtered.
_TOOL_TABLE: dict[str, str] = {
    "search_companies": "公司",
    "get_company": "公司",
    "search_assets": "资产",
    "get_asset": "资产",
    "search_clinical": "临床",
    "search_deals": "交易",
}


def _strip_tool_result(name: str, result, can_see_internal: bool):
    """Apply field_policy to a tool's return value. External users only.

    Handles the 6 individual CRM tools plus search_global (which returns a
    dict of {companies, assets, deals, clinical} sub-lists).
    """
    if can_see_internal or result is None:
        return result

    table = _TOOL_TABLE.get(name)
    if table:
        # strip_hidden handles both dict (single row) and list[dict]
        if isinstance(result, dict) and "error" in result:
            return result
        return strip_hidden(result, table, False)

    if name == "search_global" and isinstance(result, dict):
        pairs = (("companies", "公司"), ("assets", "资产"),
                 ("deals", "交易"), ("clinical", "临床"))
        for key, tbl in pairs:
            if key in result and result[key]:
                result[key] = strip_hidden(result[key], tbl, False)
        return result

    return result


def _execute_tool(
    name: str,
    inp: dict,
    user_id: str | None = None,
    can_see_internal: bool = False,
) -> str:
    """Execute a tool and return JSON-serialized result (truncated for context).

    `can_see_internal` = True for admins and internal-flagged users.
    When False, hidden CRM fields (BD priority, Q scores, strategic analysis,
    internal notes, etc.) are stripped from the returned rows before the LLM
    ever sees them — matches the REST-side field_policy enforcement.
    """
    fn = TOOL_IMPL.get(name)
    if not fn:
        return json.dumps({"error": f"Unknown tool: {name}"})

    # Block count_by on hidden columns for external users — otherwise they
    # could extract the distribution of internal fields (e.g. BD priority counts).
    if name == "count_by" and not can_see_internal:
        table = (inp or {}).get("table", "")
        group_by = (inp or {}).get("group_by", "")
        if group_by in HIDDEN_FIELDS.get(table, set()):
            return json.dumps({"error": f"字段 '{group_by}' 不对外部用户开放"}, ensure_ascii=False)

    try:
        kwargs = dict(inp or {})
        if user_id and name in ("add_to_watchlist",):
            kwargs["_user_id"] = user_id
        elif user_id and name.startswith("generate_"):
            kwargs["_user_id"] = user_id
        result = fn(**kwargs)

        # Enforce three-tier field visibility before handing data to the LLM.
        result = _strip_tool_result(name, result, can_see_internal)

        s = json.dumps(result, ensure_ascii=False, default=str)
        if len(s) > 8000:
            s = s[:8000] + "\n...[truncated]"
        return s
    except Exception as e:
        logger.exception("Tool %s failed", name)
        return json.dumps({"error": str(e)})


# ── Context entity extraction ────────────────────────────────

import re as _re
from urllib.parse import quote as _quote


def _slugify(text: str) -> str:
    """Lowercase, strip punct/whitespace for dedup keys."""
    if not text:
        return ""
    t = _re.sub(r"\s+", "_", str(text).strip().lower())
    t = _re.sub(r"[^\w\u4e00-\u9fff_\-]", "", t)
    return t[:80]


def _format_field(label: str, value) -> dict | None:
    if value is None or value == "" or str(value).strip() in ("-", "None", "null"):
        return None
    s = str(value).strip()
    if len(s) > 60:
        s = s[:57] + "..."
    return {"label": label, "value": s}


def _entity_from_company(row: dict) -> dict | None:
    name = row.get("客户名称")
    if not name:
        return None
    fields = [
        _format_field("Country", row.get("所处国家")),
        _format_field("Type", row.get("客户类型")),
        _format_field("Stage", row.get("核心产品的阶段")),
        _format_field("Disease", row.get("疾病领域")),
        _format_field("BD Priority", row.get("BD跟进优先级")),
    ]
    fields = [f for f in fields if f]
    subtitle_parts = [row.get("所处国家"), row.get("客户类型")]
    subtitle = " · ".join([p for p in subtitle_parts if p])
    return {
        "type": "context_entity",
        "id": f"company:{_slugify(name)}",
        "entity_type": "company",
        "title": str(name),
        "subtitle": subtitle or None,
        "fields": fields,
        "href": f"/companies/{_quote(str(name), safe='')}",
    }


def _entity_from_asset(row: dict) -> dict | None:
    name = row.get("资产名称") or row.get("资产代号")
    company = row.get("所属客户")
    if not name:
        return None
    fields = [
        _format_field("Company", company),
        _format_field("Phase", row.get("临床阶段")),
        _format_field("Target", row.get("靶点")),
        _format_field("MOA", row.get("作用机制(MOA)")),
        _format_field("Q Score", row.get("Q总分")),
    ]
    fields = [f for f in fields if f]
    subtitle = company or row.get("疾病领域") or ""
    href = None
    if company and name:
        href = f"/assets/{_quote(str(company), safe='')}/{_quote(str(name), safe='')}"
    return {
        "type": "context_entity",
        "id": f"asset:{_slugify(f'{company}_{name}')}",
        "entity_type": "asset",
        "title": str(name),
        "subtitle": subtitle or None,
        "fields": fields,
        "href": href,
    }


def _entity_from_clinical(row: dict) -> dict | None:
    tid = row.get("试验ID") or row.get("记录ID")
    if not tid:
        return None
    fields = [
        _format_field("Company", row.get("公司名称")),
        _format_field("Asset", row.get("资产名称")),
        _format_field("Phase", row.get("临床期次")),
        _format_field("Indication", row.get("适应症")),
        _format_field("Result", row.get("结果判定")),
    ]
    fields = [f for f in fields if f]
    record_id = row.get("记录ID") or tid
    return {
        "type": "context_entity",
        "id": f"clinical:{_slugify(str(record_id))}",
        "entity_type": "clinical",
        "title": str(tid),
        "subtitle": row.get("公司名称") or None,
        "fields": fields,
        "href": f"/clinical/{_quote(str(record_id), safe='')}",
    }


def _entity_from_deal(row: dict) -> dict | None:
    name = row.get("交易名称")
    if not name:
        return None
    fields = [
        _format_field("Type", row.get("交易类型")),
        _format_field("Buyer", row.get("买方公司")),
        _format_field("Seller", row.get("卖方/合作方")),
        _format_field("Upfront", row.get("首付款($M)") and f"${row.get('首付款($M)')}M"),
        _format_field("Total", row.get("交易总额($M)") and f"${row.get('交易总额($M)')}M"),
        _format_field("Date", row.get("宣布日期")),
    ]
    fields = [f for f in fields if f]
    return {
        "type": "context_entity",
        "id": f"deal:{_slugify(str(name))}",
        "entity_type": "deal",
        "title": str(name),
        "subtitle": row.get("交易类型") or None,
        "fields": fields,
        "href": f"/deals/{_quote(str(name), safe='')}",
    }


def _entity_from_patent(row: dict) -> dict | None:
    pid = row.get("专利号")
    if not pid:
        return None
    fields = [
        _format_field("Holder", row.get("专利持有人")),
        _format_field("Company", row.get("关联公司")),
        _format_field("Asset", row.get("关联资产")),
        _format_field("Expiry", row.get("到期日")),
        _format_field("Status", row.get("状态")),
        _format_field("Region", row.get("管辖区")),
    ]
    fields = [f for f in fields if f]
    return {
        "type": "context_entity",
        "id": f"patent:{_slugify(str(pid))}",
        "entity_type": "patent",
        "title": str(pid),
        "subtitle": row.get("专利持有人") or None,
        "fields": fields,
        "href": f"/ip/{_quote(str(pid), safe='')}",
    }


def _entity_from_buyer(row: dict) -> dict | None:
    name = row.get("company_name")
    if not name:
        return None
    fields = [
        _format_field("Focus TA", row.get("focus_therapeutic_areas")),
        _format_field("Recent Deals", row.get("recent_deal_count")),
        _format_field("Strategy", row.get("bd_strategy")),
    ]
    fields = [f for f in fields if f]
    return {
        "type": "context_entity",
        "id": f"buyer:{_slugify(str(name))}",
        "entity_type": "buyer",
        "title": str(name),
        "subtitle": "MNC Buyer Profile",
        "fields": fields,
        "href": f"/buyers/{_quote(str(name), safe='')}",
    }


def _extract_context_entities(tool_name: str, raw_result_str: str) -> list[dict]:
    """Parse a tool's JSON result string and extract structured entities for the context panel."""
    if tool_name in ("count_by", "search_global"):
        return []
    try:
        data = json.loads(raw_result_str)
    except (json.JSONDecodeError, TypeError):
        return []

    # Single-row tools return a dict; search_* tools return a list
    rows: list[dict] = []
    if isinstance(data, list):
        rows = [r for r in data if isinstance(r, dict)]
    elif isinstance(data, dict):
        rows = [data]

    if not rows:
        return []

    # Cap at top-3 entities per tool call
    rows = rows[:3]

    extractors = {
        "search_companies": _entity_from_company,
        "get_company": _entity_from_company,
        "search_assets": _entity_from_asset,
        "get_asset": _entity_from_asset,
        "search_clinical": _entity_from_clinical,
        "search_deals": _entity_from_deal,
        "search_patents": _entity_from_patent,
        "get_buyer_profile": _entity_from_buyer,
    }

    fn = extractors.get(tool_name)
    if not fn:
        return []

    entities: list[dict] = []
    for r in rows:
        try:
            e = fn(r)
            if e:
                entities.append(e)
        except Exception as ex:
            logger.warning("Failed to extract entity from %s: %s", tool_name, ex)
    return entities


# ── Attachment text extraction ───────────────────────────────

def _extract_pdf_text(filepath: Path) -> str:
    """Extract text from PDF with PyMuPDF.

    Strategy per page:
    1. Try the embedded text layer (instant, perfect for digital PDFs).
    2. If the page yields < 50 chars it is almost certainly a scanned image —
       fall back to Tesseract OCR via PyMuPDF's built-in bridge (chi_sim+eng).
       Tesseract is optional — if not installed, scanned pages are skipped
       and only text-layer pages are returned.

    Caps at 20 pages and 30 000 chars to keep prompt budgets sane.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.error("PyMuPDF not installed. Run: pip install PyMuPDF")
        return ""

    doc = fitz.open(str(filepath))
    page_texts: list[str] = []
    MAX_PAGES = 20
    ocr_available: bool | None = None   # lazily detected

    for page_num, page in enumerate(doc):
        if page_num >= MAX_PAGES:
            break

        text = page.get_text("text").strip()

        if len(text) < 50:          # sparse / image-only page — try OCR
            if ocr_available is False:
                pass  # already know tesseract is absent — skip silently
            else:
                try:
                    tp = page.get_textpage_ocr(language="chi_sim+eng", dpi=150, full=False)
                    text = page.get_text(textpage=tp).strip()
                    ocr_available = True
                except Exception as ocr_err:
                    err_str = str(ocr_err)
                    if "tesseract" in err_str.lower() or "not found" in err_str.lower():
                        if ocr_available is None:
                            logger.warning(
                                "Tesseract not found — OCR disabled. "
                                "Install: sudo apt-get install tesseract-ocr tesseract-ocr-chi-sim"
                            )
                        ocr_available = False
                    else:
                        logger.debug("OCR failed p%d of %s: %s", page_num + 1, filepath.name, ocr_err)

        if text:
            page_texts.append(f"[Page {page_num + 1}]\n{text}")

    doc.close()
    combined = "\n\n".join(page_texts)
    if not combined.strip():
        logger.warning("PDF extraction yielded no text for %s (pages=%d, ocr_available=%s)",
                       filepath.name, len(doc), ocr_available)
    else:
        logger.info("PDF extracted %d chars from %s (ocr_available=%s)", len(combined), filepath.name, ocr_available)
    return combined[:30000]


def _extract_text(filename: str) -> str:
    """Extract text from PDF/PPTX/DOCX files in BP_DIR. Returns empty string on failure."""
    filepath = BP_DIR / Path(filename).name
    if not filepath.exists():
        logger.warning("Attachment not found at %s", filepath)
        return ""
    ext = filepath.suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf_text(filepath)
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(str(filepath))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)[:20000]
        elif ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(str(filepath))
            return "\n".join(p.text for p in doc.paragraphs)[:20000]
    except Exception as e:
        logger.warning("Failed to extract %s: %s", filename, e)
        return ""
    return ""


# ── Chat streaming ───────────────────────────────────────────

class PlanConfirm(BaseModel):
    plan_id: str
    plan_title: str
    selected_steps: list[dict]      # [{id, title, description, tools_expected}]
    original_message: str


class ChatRequest(BaseModel):
    message: str
    session_id: str = "default"
    file_ids: list[str] = []
    model_id: str | None = None   # selected from /api/models; falls back to default
    # Plan mode: "auto" (heuristic), "on" (always plan), "off" (skip planning)
    plan_mode: str = "auto"
    plan_confirm: PlanConfirm | None = None
    # user_id / is_admin / is_internal are injected by the endpoint, not sent by the client
    user_id: str | None = None
    is_admin: bool = False
    is_internal: bool = False


async def _call_minimax_stream(
    client: httpx.AsyncClient,
    messages: list,
    tool_results_buffer: list,
    model: ModelSpec,
    usage_accum: dict,
    system_prompt: str | None = None,
):
    """Single call to the selected LLM; yields SSE events and collects tool_use blocks.

    `usage_accum` is a dict with keys {input_tokens, output_tokens} that this
    function adds to as the provider reports usage in `message_start` /
    `message_delta` events (Anthropic-compat shape).

    `system_prompt` overrides the default SYSTEM_PROMPT — used to inject
    plan constraints when executing an approved plan.
    """
    body = {
        "model": model.api_model,
        "system": system_prompt or SYSTEM_PROMPT,
        "messages": messages,
        "max_tokens": 4096,
        "stream": True,
        "tools": TOOLS,
    }
    headers = {
        "x-api-key": model.api_key,
        "Content-Type": "application/json",
    }
    if model.anthropic_version:
        headers["anthropic-version"] = model.anthropic_version

    collected_content = []
    current_tool_use = None
    current_text = ""
    stop_reason = None

    # 529 = server overloaded — retry up to 3 times with backoff
    max_retries = 3
    for attempt in range(max_retries):
        async with client.stream("POST", model.api_url, json=body, headers=headers) as resp:
            if resp.status_code == 529:
                await resp.aread()  # drain body
                if attempt < max_retries - 1:
                    wait = 2 ** attempt  # 1s, 2s
                    logger.warning("MiniMax 529 (attempt %d/%d), retry in %ds", attempt + 1, max_retries, wait)
                    await asyncio.sleep(wait)
                    continue
                yield ("error", "AI服务当前负载较高，请稍等片刻后重试。")
                return
            if resp.status_code != 200:
                err = await resp.aread()
                yield ("error", f"AI服务异常（{resp.status_code}），请稍后重试。")
                return

            # ── success: process the stream ──────────────────────────
            buffer = ""
            async for raw_chunk in resp.aiter_bytes():
                buffer += raw_chunk.decode("utf-8", errors="replace")
                while "\n" in buffer:
                    line, buffer = buffer.split("\n", 1)
                    line = line.strip()
                    if not line.startswith("data: "):
                        continue
                    try:
                        data = json.loads(line[6:])
                    except json.JSONDecodeError:
                        continue

                    et = data.get("type", "")

                    # ── Usage tracking (Anthropic-compat shape) ──
                    # message_start carries the initial {input_tokens, output_tokens=N}
                    # message_delta.usage updates the running output_tokens total
                    if et == "message_start":
                        u = (data.get("message", {}) or {}).get("usage") or {}
                        usage_accum["input_tokens"] += int(u.get("input_tokens") or 0)
                        usage_accum["output_tokens"] += int(u.get("output_tokens") or 0)
                    elif et == "message_delta":
                        u = data.get("usage") or {}
                        # Providers report the *final* output_tokens here, not a delta.
                        # Stash it and reconcile when we see message_stop.
                        if "output_tokens" in u:
                            usage_accum["_pending_output"] = int(u["output_tokens"] or 0)

                    if et == "content_block_start":
                        block = data.get("content_block", {})
                        if block.get("type") == "tool_use":
                            current_tool_use = {
                                "id": block.get("id"),
                                "name": block.get("name"),
                                "input_json": "",
                            }
                            yield ("tool_call_start", {"name": block.get("name")})
                        elif block.get("type") == "text":
                            current_text = ""

                    elif et == "content_block_delta":
                        delta = data.get("delta", {})
                        dt = delta.get("type", "")
                        if dt == "text_delta":
                            text = delta.get("text", "")
                            current_text += text
                            yield ("chunk", text)
                        elif dt == "input_json_delta" and current_tool_use is not None:
                            current_tool_use["input_json"] += delta.get("partial_json", "")

                    elif et == "content_block_stop":
                        if current_tool_use is not None:
                            try:
                                inp = json.loads(current_tool_use["input_json"] or "{}")
                            except json.JSONDecodeError:
                                inp = {}
                            collected_content.append({
                                "type": "tool_use",
                                "id": current_tool_use["id"],
                                "name": current_tool_use["name"],
                                "input": inp,
                            })
                            tool_results_buffer.append(current_tool_use)
                            current_tool_use = None
                        elif current_text:
                            collected_content.append({"type": "text", "text": current_text})
                            current_text = ""

                    elif et == "message_delta":
                        stop_reason = data.get("delta", {}).get("stop_reason")

                    elif et == "message_stop":
                        # Promote the pending output_tokens total from message_delta
                        # into the accumulator, then reset for the next round.
                        pending = usage_accum.pop("_pending_output", 0)
                        if pending:
                            usage_accum["output_tokens"] += pending
                        break

        # Successful stream completed — yield end event and stop retrying
        yield ("_end", {"stop_reason": stop_reason, "content": collected_content})
        return


async def _stream_chat(req: ChatRequest):
    """Main chat handler: tool-use loop with streaming."""
    session_id = req.session_id
    user_id = req.user_id
    # Field visibility: admins + internal employees see everything;
    # external users get HIDDEN_FIELDS stripped from tool results.
    can_see_internal = bool(req.is_admin or req.is_internal)

    # Resolve model (falls back to default if unknown/None)
    model = resolve_model(req.model_id)

    # If the client confirmed a plan, inject its constraints into the system
    # prompt so the LLM sticks to the approved steps.
    active_system_prompt: str | None = None
    if req.plan_confirm and req.plan_confirm.selected_steps:
        active_system_prompt = SYSTEM_PROMPT + planner_mod.build_plan_constraint(
            req.plan_confirm.plan_title,
            req.plan_confirm.selected_steps,
        )

    # Usage accumulator — incremented inside _call_minimax_stream across
    # all tool-use rounds in this request. Debited from the user's balance
    # once the stream completes.
    usage_accum = {"input_tokens": 0, "output_tokens": 0}

    # Ensure session exists in Postgres (auto-create if needed)
    if user_id:
        _ensure_session(session_id, user_id)

    # Load history from Postgres instead of in-memory dict
    history = _load_history(session_id)

    # When executing a confirmed plan, the user message + empty plan-placeholder
    # assistant message are already in history from the planner phase. Strip
    # the empty assistant turn so the LLM sees the user prompt as the latest
    # turn (and we skip re-appending + re-saving below).
    is_executing_plan = req.plan_confirm is not None
    if is_executing_plan:
        while history and history[-1].get("role") == "assistant":
            c = history[-1].get("content")
            empty = (
                not c
                or (isinstance(c, str) and not c.strip())
                or (
                    isinstance(c, list)
                    and not any(
                        (isinstance(b, dict) and b.get("text", "").strip())
                        for b in c
                    )
                )
            )
            if empty:
                history.pop()
            else:
                break

    user_text = req.message
    attachments_json = None
    if req.file_ids:
        attachment_parts = []
        failed_extractions = []
        for fid in req.file_ids:
            # Run blocking I/O (disk read + Tesseract OCR) off the event loop
            extracted = await asyncio.to_thread(_extract_text, fid)
            if extracted:
                logger.info("Extracted %d chars from attachment: %s", len(extracted), fid)
                attachment_parts.append(f"\n\n[附件内容: {fid}]\n{extracted}")
            else:
                logger.warning("Failed to extract content from attachment: %s", fid)
                failed_extractions.append(fid)

        if attachment_parts:
            user_text = req.message + "".join(attachment_parts)
            user_text += (
                "\n\n[系统指令：用户上传了文件，文件内容已在上方提供。请立即执行以下步骤："
                "1) 从文件内容中提取公司名、资产名、靶点、适应症等关键实体；"
                "2) 用提取的实体并行调用 search_companies、search_assets、search_clinical、search_deals 查询CRM数据；"
                "3) 如果文件涉及特定疾病，调用 query_treatment_guidelines 查询治疗格局；"
                "4) 综合文件内容和CRM数据，输出结构化分析报告（公司画像、管线评估、治疗格局、交易参考、BD建议）。"
                "不要只总结文件内容，必须交叉验证CRM数据。]"
            )
        elif failed_extractions:
            # Extraction failed for all files — tell the AI so it can explain honestly
            user_text = req.message + (
                f"\n\n[系统提示：用户上传了文件 {', '.join(failed_extractions)}，"
                "但文件内容无法提取（可能是加密PDF或格式不支持）。"
                "请直接告知用户文件无法解析，并询问他们能否提供文字版或核心信息。]"
            )
        attachments_json = json.dumps(req.file_ids)

    if is_executing_plan:
        # User message was already saved + already in history from plan phase.
        pass
    else:
        history.append({"role": "user", "content": user_text})

        # Persist the user message
        _save_message(session_id, "user", user_text, attachments_json=attachments_json)

    # Auto-compact: strip old tool blocks; summarize if still over budget.
    # The final history has at most ~COMPACT_TOKEN_BUDGET tokens regardless
    # of how long the session gets.
    history = await _compact_if_needed(session_id, history, model)

    # Collect entities across all tool-call rounds for persistence
    all_entities: list[dict] = []

    try:
        async with httpx.AsyncClient(timeout=300) as client:
            for _iteration in range(8):  # Up to 8 tool-call rounds for complex research
                tool_results_buffer = []
                final_content = None
                final_stop_reason = None

                async for event_type, payload in _call_minimax_stream(
                    client, history, tool_results_buffer, model, usage_accum,
                    system_prompt=active_system_prompt,
                ):
                    if event_type == "chunk":
                        yield f"data: {json.dumps({'type': 'chunk', 'content': payload})}\n\n"
                    elif event_type == "tool_call_start":
                        yield f"data: {json.dumps({'type': 'tool_call', 'name': payload['name']})}\n\n"
                    elif event_type == "error":
                        yield f"data: {json.dumps({'type': 'error', 'message': payload})}\n\n"
                        yield f"data: {json.dumps({'type': 'done'})}\n\n"
                        return
                    elif event_type == "_end":
                        final_content = payload["content"]
                        final_stop_reason = payload["stop_reason"]

                if final_content is None:
                    break

                history.append({"role": "assistant", "content": final_content})

                if final_stop_reason == "tool_use" and tool_results_buffer:
                    tool_results_msg = []
                    tool_events = []
                    for tu in tool_results_buffer:
                        try:
                            inp = json.loads(tu["input_json"] or "{}")
                        except json.JSONDecodeError:
                            inp = {}
                        result_str = _execute_tool(
                            tu["name"], inp,
                            user_id=user_id,
                            can_see_internal=can_see_internal,
                        )
                        yield f"data: {json.dumps({'type': 'tool_result', 'name': tu['name']})}\n\n"

                        # Emit report_task event for async report tools so the
                        # frontend can show an inline polling card.
                        if tu["name"].startswith("generate_"):
                            try:
                                _rd = json.loads(result_str)
                                if _rd.get("status") == "queued" and _rd.get("task_id"):
                                    yield f"data: {json.dumps({'type': 'report_task', 'task_id': _rd['task_id'], 'slug': tu['name'].replace('generate_', '', 1), 'estimated_seconds': _rd.get('estimated_seconds', 60)})}\n\n"
                            except Exception:
                                pass

                        # Collect tool event for persistence
                        tool_events.append({
                            "name": tu["name"],
                            "input": inp,
                            "result_preview": result_str[:500],
                        })

                        # Emit context_entity events for each extracted entity
                        for entity in _extract_context_entities(tu["name"], result_str):
                            yield f"data: {json.dumps(entity, ensure_ascii=False)}\n\n"
                            all_entities.append(entity)

                        tool_results_msg.append({
                            "type": "tool_result",
                            "tool_use_id": tu["id"],
                            "content": result_str,
                        })

                    # Persist the assistant message (with tool calls) and tool results
                    tools_json = json.dumps(tool_events, ensure_ascii=False, default=str)
                    _save_message(session_id, "assistant", final_content, tools_json=tools_json)
                    _save_message(session_id, "user", tool_results_msg)

                    history.append({"role": "user", "content": tool_results_msg})
                    continue
                else:
                    # Final assistant text (no more tool calls) — persist it
                    _save_message(session_id, "assistant", final_content)
                    break

        # Persist all extracted context entities
        _save_entities(session_id, all_entities)

        # ── Bill credits for this request ──
        # Debit after the stream finishes so we never charge for a failed turn.
        # Admin users are never billed.
        credits_charged = 0.0
        balance_remaining = None
        if user_id and not req.is_admin:
            credits_charged = credits_mod.record_usage(
                user_id=user_id,
                session_id=session_id,
                model_id=model.id,
                input_tokens=usage_accum.get("input_tokens", 0),
                output_tokens=usage_accum.get("output_tokens", 0),
                input_weight=model.input_weight,
                output_weight=model.output_weight,
            )
            try:
                balance_info = credits_mod.get_balance(user_id)
                balance_remaining = balance_info["balance"]
            except Exception:
                balance_remaining = None

        yield (
            "data: "
            + json.dumps({
                "type": "usage",
                "model": model.id,
                "input_tokens": usage_accum.get("input_tokens", 0),
                "output_tokens": usage_accum.get("output_tokens", 0),
                "credits_charged": credits_charged,
                "balance": balance_remaining,
            })
            + "\n\n"
        )
        yield f"data: {json.dumps({'type': 'done'})}\n\n"

    except httpx.TimeoutException:
        yield f"data: {json.dumps({'type': 'error', 'message': 'Request timed out'})}\n\n"
        yield f"data: {json.dumps({'type': 'done'})}\n\n"
    except Exception as e:
        logger.exception("Chat error")
        yield f"data: {json.dumps({'type': 'error', 'message': str(e)[:500]})}\n\n"
        yield f"data: {json.dumps({'type': 'done'})}\n\n"


_PLAN_KEYWORDS = (
    "分析", "深度", "画像", "scout", "找出", "筛选",
    "BD机会", "BD 机会", "机会分析", "对比", "评估",
    "深入", "盘点", "梳理", "报告", "全面",
)


def _should_plan(message: str) -> bool:
    """Heuristic: decide if an auto-mode request warrants a plan phase."""
    if not message or len(message.strip()) < 8:
        return False
    if len(message) > 150:
        return True
    return any(kw in message for kw in _PLAN_KEYWORDS)


async def _stream_plan_only(req: ChatRequest):
    """Phase 1 stream: generate a plan proposal, emit one SSE event, done.

    Falls back to normal execution (_stream_chat) if planner fails.
    """
    model = resolve_model(req.model_id)
    # Load short history for context
    history = _load_history(req.session_id)

    # Ensure the session exists, but DON'T save the user message yet — we
    # only save it after we know whether planner succeeds. This prevents
    # double-save if planner fails and we fall back to _stream_chat (which
    # saves the user message itself).
    if req.user_id:
        _ensure_session(req.session_id, req.user_id)

    plan = await planner_mod.generate_plan(req.message, history, model)

    if plan is None:
        # Planner failed — fall through to normal execution. _stream_chat
        # will save the user message as part of its normal flow.
        logger.warning("Planner returned no plan; falling back to normal execution")
        fallback_req = req.model_copy(update={"plan_mode": "off"})
        async for chunk in _stream_chat(fallback_req):
            yield chunk
        return

    # Planner succeeded — now we commit to the plan flow, save user msg.
    if req.user_id:
        _save_message(req.session_id, "user", req.message)

    # Bill planner tokens (admins skip)
    usage = plan.pop("_usage", {}) or {}
    credits_charged = 0.0
    balance_remaining = None
    if req.user_id and not req.is_admin:
        try:
            credits_charged = credits_mod.record_usage(
                user_id=req.user_id,
                session_id=req.session_id,
                model_id=model.id,
                input_tokens=usage.get("input_tokens", 0),
                output_tokens=usage.get("output_tokens", 0),
                input_weight=model.input_weight,
                output_weight=model.output_weight,
            )
            balance_remaining = credits_mod.get_balance(req.user_id)["balance"]
        except Exception:
            logger.exception("Failed to record planner usage")

    # Persist the plan as a placeholder assistant message so it survives reload
    if req.user_id:
        _save_message(
            req.session_id,
            "assistant",
            "",
            tools_json=json.dumps({"plan": plan, "original_message": req.message}, ensure_ascii=False),
        )

    yield f"data: {json.dumps({'type': 'plan_proposal', 'plan': plan, 'original_message': req.message}, ensure_ascii=False)}\n\n"
    if credits_charged:
        yield f"data: {json.dumps({'type': 'usage', 'credits_charged': credits_charged, 'balance': balance_remaining})}\n\n"
    yield f"data: {json.dumps({'type': 'done'})}\n\n"


@router.post("")
async def chat_stream(req: ChatRequest, user: dict = Depends(get_current_user)):
    req.user_id = user["id"]
    req.is_admin = bool(user.get("is_admin"))
    req.is_internal = bool(user.get("is_internal"))
    # Admin users bypass credit check (still logged, just never blocked).
    if not req.is_admin:
        # Fail fast with a clean 402 if the user is out of credits — never mid-stream.
        credits_mod.ensure_balance(user["id"])

    # Decide: planning phase, execution phase, or normal (no-plan) flow.
    is_confirming = req.plan_confirm is not None
    should_plan_now = (
        not is_confirming
        and req.plan_mode != "off"
        and (req.plan_mode == "on" or _should_plan(req.message))
    )

    if should_plan_now:
        generator = _stream_plan_only(req)
    else:
        # Normal execution — either no plan mode, or user confirmed a plan.
        # When confirming, use the original message as the chat input so the
        # LLM sees the actual question (not "executed plan confirmation").
        if is_confirming and req.plan_confirm:
            req = req.model_copy(update={"message": req.plan_confirm.original_message})
        generator = _stream_chat(req)

    return StreamingResponse(
        generator,
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive", "X-Accel-Buffering": "no"},
    )
