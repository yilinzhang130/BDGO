"""Attachment text extraction for PDF / PPTX / DOCX uploads.

Used when the user attaches a file to a chat turn — extracted text is
appended to the user message so the LLM can reason about it alongside
CRM data.

Caps are intentional: 20 pages or 30_000 chars for PDFs, 20_000 for
slide decks and Word docs. Prevents one huge deck from drowning the
context budget.
"""

from __future__ import annotations

import logging
from pathlib import Path

from config import BP_DIR

logger = logging.getLogger(__name__)


def extract_pdf_text(filepath: Path) -> str:
    """Extract text from a PDF with PyMuPDF.

    Strategy per page:
      1. Try the embedded text layer (instant, perfect for digital PDFs).
      2. If the page yields < 50 chars it is almost certainly a scanned
         image — fall back to Tesseract OCR via PyMuPDF's built-in bridge
         (chi_sim+eng). Tesseract is optional — if not installed, scanned
         pages are skipped and only text-layer pages are returned.

    Caps at 20 pages and 30_000 chars to keep prompt budgets sane.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.error("PyMuPDF not installed. Run: pip install PyMuPDF")
        return ""

    doc = fitz.open(str(filepath))
    try:
        page_texts: list[str] = []
        MAX_PAGES = 20
        ocr_available: bool | None = None   # lazily detected

        for page_num, page in enumerate(doc):
            if page_num >= MAX_PAGES:
                break

            text = page.get_text("text").strip()

            if len(text) < 50:          # sparse / image-only page — try OCR
                if ocr_available is False:
                    pass  # already know tesseract is absent — skip silently
                else:
                    try:
                        tp = page.get_textpage_ocr(language="chi_sim+eng", dpi=150, full=False)
                        text = page.get_text(textpage=tp).strip()
                        ocr_available = True
                    except Exception as ocr_err:
                        err_str = str(ocr_err)
                        if "tesseract" in err_str.lower() or "not found" in err_str.lower():
                            if ocr_available is None:
                                logger.warning(
                                    "Tesseract not found — OCR disabled. "
                                    "Install: sudo apt-get install tesseract-ocr tesseract-ocr-chi-sim"
                                )
                            ocr_available = False
                        else:
                            logger.debug("OCR failed p%d of %s: %s",
                                         page_num + 1, filepath.name, ocr_err)

            if text:
                page_texts.append(f"[Page {page_num + 1}]\n{text}")

        combined = "\n\n".join(page_texts)
        if not combined.strip():
            logger.warning(
                "PDF extraction yielded no text for %s (pages=%d, ocr_available=%s)",
                filepath.name, len(doc), ocr_available,
            )
        else:
            logger.info(
                "PDF extracted %d chars from %s (ocr_available=%s)",
                len(combined), filepath.name, ocr_available,
            )
        return combined[:30_000]
    finally:
        doc.close()


def extract_text(filename: str) -> str:
    """Extract text from PDF / PPTX / DOCX files in BP_DIR.

    Returns empty string on failure (caller logs).
    """
    filepath = BP_DIR / Path(filename).name
    if not filepath.exists():
        logger.warning("Attachment not found at %s", filepath)
        return ""

    ext = filepath.suffix.lower()
    try:
        if ext == ".pdf":
            return extract_pdf_text(filepath)
        if ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs = Presentation(str(filepath))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
            return "\n".join(parts)[:20_000]
        if ext in (".docx", ".doc"):
            from docx import Document
            doc = Document(str(filepath))
            return "\n".join(p.text for p in doc.paragraphs)[:20_000]
    except Exception as e:
        logger.warning("Failed to extract %s: %s", filename, e)
        return ""

    return ""
