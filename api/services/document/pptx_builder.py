"""Lightweight pptx builder for Deal Teaser reports.

Produces an 8-page BD-styled deck. Consumes a structured dict produced by
the LLM — one slide per section. Each slide is built from a thin helper
that places a title bar + body text / table.
"""

from __future__ import annotations

import io
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x4E, 0x79)
GOLD = RGBColor(0xBF, 0x92, 0x3C)
DARK_GRAY = RGBColor(0x40, 0x40, 0x40)
LIGHT_GRAY = RGBColor(0xE7, 0xEB, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x12, 0x1A, 0x24)


@dataclass
class TeaserContent:
    """Structured deck content — filled by the LLM, rendered by build_deck()."""

    asset_name: str
    indication: str
    company: str
    target: str
    moa: str
    phase: str
    date: str
    highlights: list[str]  # 3-5 bullets with embedded data
    unmet_need: str  # ~200 字
    tam: str  # "$8.5B global"
    peak_revenue: str  # "$1.2-1.8B"
    mechanism: str  # MoA + differentiation vs SoC
    clinical_data: list[list[str]]  # table rows: ["Metric", "Our Asset", "SoC", "Δ"]
    competitive: list[list[str]]  # ["Competitor", "Company", "Stage", "Differentiation"]
    development_plan: list[str]  # milestone bullets
    deal_structure: str
    contact: str


def _add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.9)
    )
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    tf = bar.text_frame
    tf.margin_left = Inches(0.35)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    for r in p.runs:
        r.font.size = Pt(22)
        r.font.bold = True
        r.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        for r in p2.runs:
            r.font.size = Pt(12)
            r.font.color.rgb = GOLD

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.9), Inches(13.33), Inches(0.05)
    )
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD


def _add_footer(slide, text: str = "CONFIDENTIAL — For Discussion Purposes Only") -> None:
    box = slide.shapes.add_textbox(Inches(0.35), Inches(7.0), Inches(12.6), Inches(0.3))
    tf = box.text_frame
    tf.margin_left = Inches(0)
    p = tf.paragraphs[0]
    p.text = text
    for r in p.runs:
        r.font.size = Pt(9)
        r.font.italic = True
        r.font.color.rgb = DARK_GRAY


def _add_body_text(
    slide, text: str, *, top: float = 1.2, height: float = 5.5, size: int = 14
) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split("\n"):
        line = line.rstrip()
        if not line:
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = BLACK


def _add_bullets(slide, bullets: list[str], *, top: float = 1.2, size: int = 16) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.space_after = Pt(12)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = BLACK


def _add_table(slide, headers: list[str], rows: list[list[str]], *, top: float = 1.2) -> None:
    if not rows:
        _add_body_text(slide, "(no data)", top=top)
        return
    cols = len(headers)
    nrows = len(rows) + 1
    left = Inches(0.6)
    width = Inches(12.1)
    height = Inches(min(5.5, 0.5 + 0.5 * nrows))
    tbl = slide.shapes.add_table(nrows, cols, left, Inches(top), width, height).table
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        p = cell.text_frame.paragraphs[0]
        p.text = h
        for r in p.runs:
            r.font.size = Pt(13)
            r.font.bold = True
            r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j in range(cols):
            cell = tbl.cell(i, j)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            p = cell.text_frame.paragraphs[0]
            p.text = row[j] if j < len(row) else ""
            for r in p.runs:
                r.font.size = Pt(11)
                r.font.color.rgb = BLACK


def build_deck(content: TeaserContent) -> bytes:
    """Render a pptx deck and return bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Cover ─────────────────────────────────────
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY
    title = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.2))
    tf = title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content.asset_name
    for r in p.runs:
        r.font.size = Pt(48)
        r.font.bold = True
        r.font.color.rgb = WHITE

    subtitle = s.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.9))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = f"{content.indication} · {content.phase}"
    for r in p.runs:
        r.font.size = Pt(24)
        r.font.color.rgb = GOLD

    meta = s.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.9))
    tf = meta.text_frame
    p = tf.paragraphs[0]
    p.text = f"{content.company}   |   Target: {content.target}   |   MoA: {content.moa}"
    for r in p.runs:
        r.font.size = Pt(14)
        r.font.color.rgb = WHITE

    conf = s.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.5))
    tf = conf.text_frame
    p = tf.paragraphs[0]
    p.text = f"CONFIDENTIAL — Deal Teaser · {content.date}"
    for r in p.runs:
        r.font.size = Pt(11)
        r.font.italic = True
        r.font.color.rgb = GOLD

    # ── Slide 2: Investment Highlights ─────────────────────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Investment Highlights", "Why this asset deserves buyer attention")
    _add_bullets(s, content.highlights[:5], top=1.3, size=17)
    _add_footer(s)

    # ── Slide 3: Unmet Need & Market ───────────────────────
    s = prs.slides.add_slide(blank)
    _add_title_bar(
        s, "Unmet Need & Market", f"TAM {content.tam}  ·  Peak Revenue {content.peak_revenue}"
    )
    _add_body_text(s, content.unmet_need, top=1.3)
    _add_footer(s)

    # ── Slide 4: Mechanism & Differentiation ───────────────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Mechanism & Differentiation", f"{content.target} / {content.moa}")
    _add_body_text(s, content.mechanism, top=1.3)
    _add_footer(s)

    # ── Slide 5: Clinical Data ─────────────────────────────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Clinical Data", "Key efficacy / safety readouts")
    if content.clinical_data and len(content.clinical_data) > 0:
        headers = content.clinical_data[0]
        rows = content.clinical_data[1:]
        _add_table(s, headers, rows, top=1.3)
    else:
        _add_body_text(s, "(no clinical data provided)", top=1.3)
    _add_footer(s)

    # ── Slide 6: Competitive Landscape ─────────────────────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Competitive Landscape", "Who else is in the space")
    if content.competitive and len(content.competitive) > 0:
        headers = content.competitive[0]
        rows = content.competitive[1:]
        _add_table(s, headers, rows, top=1.3)
    else:
        _add_body_text(s, "(no competitive data provided)", top=1.3)
    _add_footer(s)

    # ── Slide 7: Development Plan & Milestones ─────────────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Development Plan & Milestones", "Key catalysts and timelines")
    _add_bullets(s, content.development_plan, top=1.3, size=16)
    _add_footer(s)

    # ── Slide 8: Deal Structure & Contact ──────────────────
    s = prs.slides.add_slide(blank)
    _add_title_bar(s, "Proposed Deal Structure & Next Steps", "")
    _add_body_text(s, content.deal_structure + "\n\n" + content.contact, top=1.3)
    _add_footer(s)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
