"""Extract asset metadata from a BP file (PDF / DOCX / PPTX).

Used by the upload router to suggest a follow-up `/teaser` command after
a BP lands. Returns a best-effort dict — every field is optional and
empty strings mean "not found", not failure.

Pipeline:
  1. Plain-text extraction (PDF/DOCX via contract_extract; PPTX inline)
  2. Single LLM call with a strict JSON schema asking for company, asset,
     indication, target, phase, moa.

Failures are returned as ``{"error": "..."}`` — caller is expected to
fall through gracefully (the upload itself must always succeed).
"""

from __future__ import annotations

import json
import logging
from pathlib import Path

from services.document.contract_extract import extract_contract_text

logger = logging.getLogger(__name__)

_MAX_TEXT_CHARS = 12_000  # cap LLM input for speed/cost

# Canonical 6-field core (always extracted)
_CORE_FIELDS = ("company_name", "asset_name", "indication", "target", "phase", "moa")

# S1-03: 4 additional fields (best-effort; left empty if not found)
_EXTENDED_FIELDS = ("modality", "ip_timeline", "funding", "team")

_ASSET_FIELDS = _CORE_FIELDS + _EXTENDED_FIELDS

# Canonical modality values — the LLM must pick one or leave blank.
_MODALITY_VALUES = (
    "小分子",
    "单克隆抗体",
    "双抗/多抗",
    "ADC",
    "CAR-T",
    "TCR-T",
    "NK细胞疗法",
    "siRNA/ASO",
    "基因疗法",
    "蛋白质/多肽",
    "疫苗",
    "放射性药物",
    "其他生物制品",
)

_EXTRACT_SYSTEM = """你是 BD 文档解析助手。任务：从 BP（业务计划书 / 资产介绍）文本里抽取 10 个字段。

硬规则：
- 只输出 JSON 对象，不要加任何解释或 markdown 代码块
- 找不到的字段写空字符串 ""，不要编造
- JSON 键（共 10 个）：
    company_name  公司名称
    asset_name    资产名称（优先用代号/商品名，如 "ABC-1234"、"PD-1单抗"）
    indication    主要适应症（如 "NSCLC 一线"、"弥漫大B细胞淋巴瘤"）
    target        分子靶点（如 "PD-1"、"KRAS G12D"），不是治疗领域
    phase         开发阶段，必须是以下之一（否则留空）：
                  Preclinical / Phase 1 / Phase 2 / Phase 3 / Approved
    moa           作用机制一句话（如 "Covalent KRAS G12D inhibitor"），可留空
    modality      药物形式，必须从以下选择（否则留空）：
                  小分子 / 单克隆抗体 / 双抗/多抗 / ADC / CAR-T / TCR-T /
                  NK细胞疗法 / siRNA/ASO / 基因疗法 / 蛋白质/多肽 / 疫苗 /
                  放射性药物 / 其他生物制品
    ip_timeline   关键专利到期年份或 FTO 状态一句话（如 "化合物专利2034年到期"），
                  找不到留空
    funding       最近一轮融资（如 "Series B $80M 2024"、"A轮 5亿人民币 2023年"），
                  找不到留空
    team          1-3 位关键人物（如 "CEO: John Smith; CMO: Dr. Jane Lee"），
                  找不到留空
"""

_EXTRACT_USER = """以下是 BP 文档前 {n_chars} 字：

═══════════════════════════════
{text}
═══════════════════════════════

输出严格 JSON（不要 markdown 代码块、不要解释）：
"""


def extract_asset_metadata(filepath: Path, llm_fn) -> dict:
    """Extract asset metadata from ``filepath``.

    ``llm_fn`` is a callable matching ``ReportContext.llm`` style:
    ``llm_fn(system: str, messages: list[dict], max_tokens: int) -> str``.
    Injected so we don't take a hard dependency on the report context.
    """
    try:
        text = _extract_text(filepath)
    except Exception as e:
        logger.warning("asset_extract: text extraction failed for %s: %s", filepath, e)
        return {"error": "text_extraction_failed"}

    if not text or len(text.strip()) < 200:
        return {"error": "text_too_short"}

    text = text[:_MAX_TEXT_CHARS]
    try:
        raw = llm_fn(
            system=_EXTRACT_SYSTEM,
            messages=[
                {"role": "user", "content": _EXTRACT_USER.format(n_chars=len(text), text=text)}
            ],
            max_tokens=600,  # S1-03: bumped from 400 — 10 fields need more tokens
        )
    except Exception as e:
        logger.warning("asset_extract: LLM call failed: %s", e)
        return {"error": "llm_call_failed"}

    parsed = _parse_json_loosely(raw)
    if not parsed:
        return {"error": "json_parse_failed", "raw": raw[:200]}

    return {field: str(parsed.get(field) or "").strip() for field in _ASSET_FIELDS}


def build_teaser_command(asset: dict) -> str | None:
    """Build a slash-command string from an extracted asset dict.

    Returns ``None`` when too few required fields are present to make a
    useful suggestion. DealTeaserService requires company_name, asset_name,
    indication, target, phase — we ask for at least 3 of those.
    """
    required = ["company_name", "asset_name", "indication", "target", "phase"]
    present = [k for k in required if asset.get(k)]
    if len(present) < 3:
        return None
    parts = ["/teaser"]
    for k in required + ["moa", "modality"]:
        v = asset.get(k)
        if v:
            parts.append(f'{k}="{v}"')
    return " ".join(parts)


def build_intake_seed(asset: dict) -> str | None:
    """Build a planner-friendly seed message that triggers BD intake plan mode.

    The seed describes the asset and asks for a multi-step intake plan:
    target/disease/IP research → asset evaluation → buyer matching →
    timing → BD strategy synthesis. Phrased to (a) hit ``should_plan``
    keywords and (b) give the planner LLM enough context to pick the
    right service combination.

    Returns ``None`` if asset has too little context (need company + asset).
    """
    company = asset.get("company_name")
    asset_name = asset.get("asset_name")
    if not company or not asset_name:
        return None

    fields = []
    if asset.get("modality"):
        fields.append(f"形式 {asset['modality']}")
    if asset.get("target"):
        fields.append(f"靶点 {asset['target']}")
    if asset.get("indication"):
        fields.append(f"适应症 {asset['indication']}")
    if asset.get("phase"):
        fields.append(f"阶段 {asset['phase']}")
    if asset.get("moa"):
        fields.append(f"MoA {asset['moa']}")
    fields_str = ("（" + ", ".join(fields) + "）") if fields else ""

    # Append enrichment context so the planner LLM has richer input.
    extra_parts: list[str] = []
    if asset.get("funding"):
        extra_parts.append(f"融资背景：{asset['funding']}。")
    if asset.get("team"):
        extra_parts.append(f"核心团队：{asset['team']}。")
    if asset.get("ip_timeline"):
        extra_parts.append(f"IP 情况：{asset['ip_timeline']}。")
    extra_str = (" " + " ".join(extra_parts)) if extra_parts else ""

    return (
        f"为 {company} 的 {asset_name}{fields_str} 启动 BD intake 全面分析。"
        f"{extra_str}"
        f"请规划一个多步执行计划，覆盖：靶点深度调研、适应症 landscape、IP/FTO 检查、"
        f"资产吸引力评估、Top-3 买方匹配、接触时机建议，最后综合输出 BD 策略备忘。"
        f"每步独立可取消，让我能勾选要跑的项目。"
    )


def _extract_text(filepath: Path) -> str:
    suffix = filepath.suffix.lower()
    if suffix in (".pdf", ".docx"):
        return extract_contract_text(filepath)
    if suffix in (".pptx", ".ppt"):
        return _extract_pptx(filepath)
    if suffix == ".doc":
        # python-docx doesn't read legacy .doc; skip extraction
        return ""
    return ""


def _extract_pptx(filepath: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(filepath))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return "\n".join(parts)


def _parse_json_loosely(raw: str) -> dict | None:
    """Parse JSON from an LLM response that may have leading/trailing prose."""
    if not raw:
        return None
    raw = raw.strip()
    # Strip markdown code fences if present.
    if raw.startswith("```"):
        raw = raw.strip("`")
        first_nl = raw.find("\n")
        if first_nl != -1:
            raw = raw[first_nl + 1 :]
        if raw.endswith("```"):
            raw = raw[:-3]
    # Find first {...} balanced block.
    start = raw.find("{")
    end = raw.rfind("}")
    if start == -1 or end == -1 or end <= start:
        return None
    try:
        return json.loads(raw[start : end + 1])
    except json.JSONDecodeError:
        return None
